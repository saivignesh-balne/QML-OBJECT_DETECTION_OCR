from __future__ import annotations

import json
from datetime import date, datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
DOCS_DIR = PROJECT_ROOT / "docs"
ARTIFACTS_DIR = PROJECT_ROOT / "artifacts"
BENCHMARK_PATH = ARTIFACTS_DIR / "benchmark_report.json"
DETECTOR_SUMMARY_PATH = ARTIFACTS_DIR / "qml_ocr_detector.summary.json"
OUTPUT_PPTX = DOCS_DIR / "INTERNSHIP_EVALUATION_PRESENTATION.pptx"
FALLBACK_OUTPUT_PPTX = DOCS_DIR / "INTERNSHIP_EVALUATION_PRESENTATION_QML_FOCUSED.pptx"

FONT = "Arial"
NAVY = RGBColor(20, 36, 68)
NAVY_DARK = RGBColor(13, 24, 46)
TEAL = RGBColor(0, 128, 113)
GOLD = RGBColor(180, 140, 54)
CLASSICAL = RGBColor(34, 113, 179)
QUANTUM = RGBColor(122, 77, 196)
INK = RGBColor(32, 42, 56)
MUTED = RGBColor(88, 100, 120)
BG = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 220, 228)
SOFT_BLUE = RGBColor(234, 242, 251)
SOFT_GREEN = RGBColor(234, 247, 244)
SOFT_PURPLE = RGBColor(241, 236, 251)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def set_slide_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, section: str, index: int, total: int) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.72), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    section_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.16), Inches(6.0), Inches(0.28))
    p = section_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = section.upper()
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE

    page_box = slide.shapes.add_textbox(Inches(12.1), Inches(0.16), Inches(0.7), Inches(0.28))
    p = page_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{index}/{total}"
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = WHITE


def add_header(slide, section: str, title: str, subtitle: str, index: int, total: int) -> None:
    set_slide_background(slide)
    add_top_band(slide, section, index, total)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(11.6), Inches(0.55))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = NAVY_DARK

    subtitle_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.58), Inches(11.8), Inches(0.38))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED


def add_panel(slide, x, y, w, h, fill_color: RGBColor = WHITE, line_color: RGBColor = LINE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color


def add_panel_title(slide, x, y, w, title: str, color: RGBColor = NAVY) -> None:
    box = slide.shapes.add_textbox(x, y, w, Inches(0.25))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = color


def add_bullet_panel(
    slide,
    x,
    y,
    w,
    h,
    title: str,
    bullets: list[str],
    title_color: RGBColor = NAVY,
    fill_color: RGBColor = WHITE,
    font_size: int = 15,
) -> None:
    add_panel(slide, x, y, w, h, fill_color=fill_color)
    add_panel_title(slide, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), title, color=title_color)
    box = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.45), w - Inches(0.24), h - Inches(0.52))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(4)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(5)
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK


def add_metric_card(slide, x, y, label: str, value: str, note: str | None = None, color: RGBColor = TEAL) -> None:
    width = Inches(2.2)
    height = Inches(1.05)
    add_panel(slide, x, y, width, height)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.10), height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.color.rgb = color

    label_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), Inches(1.8), Inches(0.2))
    p = label_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED

    value_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.34), Inches(1.8), Inches(0.34))
    p = value_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(19)
    r.font.color.rgb = NAVY_DARK

    if note:
        note_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.70), Inches(1.8), Inches(0.18))
        p = note_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = note
        r.font.name = FONT
        r.font.size = Pt(8.5)
        r.font.color.rgb = MUTED


def add_text_callout(slide, x, y, w, h, title: str, body: str, fill_color: RGBColor = WHITE) -> None:
    add_panel(slide, x, y, w, h, fill_color=fill_color)
    add_panel_title(slide, x + Inches(0.15), y + Inches(0.10), w - Inches(0.3), title, color=NAVY)
    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.42), w - Inches(0.3), h - Inches(0.52))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = FONT
    p.font.size = Pt(14)
    p.font.color.rgb = INK


def format_pct(value: float) -> str:
    return f"{value * 100:.2f}%"


def short_model_label(name: str) -> str:
    replacements = {
        "Classical SVM (RBF)": "Classical SVM [C]",
        "Logistic Regression": "Logistic Regression [C]",
        "MLP Classifier": "MLP Classifier [C]",
        "Random Forest": "Random Forest [C]",
        "QSVM ZZ Kernel Linear (amplitude)": "QSVM Linear Amp [Q]",
        "QSVM Pauli Kernel (angle)": "QSVM Pauli [Q]",
        "QSVM ZZ Kernel Linear (angle)": "QSVM Linear Angle [Q]",
        "QSVM ZZ Kernel Full (angle)": "QSVM Full Angle [Q]",
        "VQC RealAmplitudes (angle)": "VQC RealAmp [Q]",
        "VQC EfficientSU2 (angle)": "VQC EfficientSU2 [Q]",
    }
    return replacements.get(name, name)


def add_horizontal_bar_chart(slide, x, y, w, h, title: str, categories: list[str], values: list[float], color: RGBColor) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series(title, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data).chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max(max(values) + 5, 100)
    chart.category_axis.tick_labels.font.size = Pt(8.5)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = color
    chart.series[0].format.line.color.rgb = color


def style_table(table, header_fill: RGBColor = NAVY, header_font: RGBColor = WHITE) -> None:
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_idx == 0 else WHITE
            cell.text_frame.word_wrap = True
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(9.5 if row_idx == 0 else 9)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = header_font if row_idx == 0 else INK


def add_model_table(slide, x, y, w, h, rows: list[dict]) -> None:
    headers = ["Model", "Family", "Accuracy", "Macro F1", "Train Time (s)"]
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for row_idx, row in enumerate(rows, start=1):
        values = [
            short_model_label(str(row.get("name", ""))),
            str(row.get("model_family", "")).title(),
            format_pct(float(row.get("accuracy", 0.0))),
            f"{float(row.get('macro_f1', 0.0)) * 100:.2f}%",
            f"{float(row.get('train_time_seconds', 0.0)):.1f}",
        ]
        for col_idx, value in enumerate(values):
            table.cell(row_idx, col_idx).text = value
    style_table(table)


def get_best_by_family(classifiers: list[dict], family: str) -> dict:
    items = [row for row in classifiers if row.get("model_family") == family]
    return max(items, key=lambda row: float(row.get("accuracy", 0.0)))


def add_title_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, color=NAVY_DARK)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.55), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(8.6), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Hybrid Quantum Machine Learning\nfor Object Detection Research and OCR"
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(3.0), Inches(7.8), Inches(1.0))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Internship Evaluation Presentation"
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(215, 223, 236)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"Research theme: Classical vs Hybrid Quantum Benchmarking | {date.today().strftime('%d %B %Y')}"
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(192, 203, 220)

    add_metric_card(slide, Inches(9.25), Inches(1.35), "Primary Objective", "QML Research", "Benchmark-driven evaluation", color=TEAL)
    add_metric_card(slide, Inches(9.25), Inches(2.65), "Benchmark Set", "3 Classes", "Bottle, chip packet, medicine box", color=GOLD)
    add_metric_card(slide, Inches(9.25), Inches(3.95), "Recommendation", "Classical Stack", "YOLO + SVM + OCR ensemble", color=CLASSICAL)
    add_metric_card(slide, Inches(9.25), Inches(5.25), "Innovation Track", "Hybrid QML", "Qiskit-based classifier comparison", color=QUANTUM)

    page_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(0.7), Inches(0.2))
    p = page_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"1/{total}"
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = WHITE


def add_agenda_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Introduction",
        "Agenda and presentation context",
        "The presentation positions the project as both a practical application and a QML research study.",
        2,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.6),
        "Agenda",
        [
            "Project background and research objective",
            "System architecture and experimental workflow",
            "Methods, technologies, and QML approach",
            "Internship outcomes and benchmark findings",
            "Recommendation, conclusion, and future scope",
        ],
        title_color=TEAL,
        fill_color=SOFT_GREEN,
        font_size=17,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.6),
        "Presentation focus",
        [
            "The project does not assume that quantum models are better by default.",
            "Instead, it evaluates classical and hybrid quantum models under the same pipeline and benchmark logic.",
            "The final recommendation is based on measured performance, training cost, and deployment practicality.",
        ],
        title_color=GOLD,
        fill_color=WHITE,
        font_size=17,
    )


def add_project_context_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Project Presentation",
        "Project context and problem definition",
        "The project is positioned as QML object-detection research and development, with three object classes used as the benchmark domain.",
        3,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(6.0),
        Inches(4.7),
        "Research context",
        [
            "The main research problem is how to evaluate QML inside a realistic object-detection and classification workflow.",
            "Real object images contain clutter, reflections, scale variation, and partially visible text.",
            "This makes the task suitable for both applied computer vision and comparative QML research.",
        ],
        title_color=NAVY,
        fill_color=WHITE,
        font_size=17,
    )
    add_bullet_panel(
        slide,
        Inches(6.95),
        Inches(2.0),
        Inches(5.7),
        Inches(4.7),
        "Project objective",
        [
            "Build a complete end-to-end pipeline for object detection, classification, and OCR.",
            "Use the pipeline as a controlled platform for classical versus hybrid quantum comparison.",
            "Treat the selected object classes as an experimental benchmark rather than the project definition itself.",
        ],
        title_color=TEAL,
        fill_color=SOFT_BLUE,
        font_size=17,
    )


def add_architecture_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Project Presentation",
        "End-to-end system architecture",
        "The architecture separates the application flow from the comparative QML evaluation layer.",
        4,
        total,
    )

    steps = [
        ("Input Image", NAVY),
        ("Preprocessing", TEAL),
        ("YOLO Detection", CLASSICAL),
        ("ROI Extraction", GOLD),
        ("Classification", QUANTUM),
        ("OCR", TEAL),
        ("Final Output", NAVY_DARK),
    ]
    x_positions = [0.65, 2.2, 3.85, 5.5, 7.15, 8.8, 10.45]
    for idx, ((label, color), xpos) in enumerate(zip(steps, x_positions)):
        add_panel(slide, Inches(xpos), Inches(2.7), Inches(1.55), Inches(1.02))
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(xpos), Inches(2.7), Inches(1.55), Inches(0.1))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.color.rgb = color
        box = slide.shapes.add_textbox(Inches(xpos + 0.05), Inches(2.95), Inches(1.45), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(12)
        r.font.color.rgb = INK
        if idx < len(steps) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(xpos + 1.55),
                Inches(3.2),
                Inches(xpos + 1.75),
                Inches(3.2),
            )
            connector.line.color.rgb = MUTED
            connector.line.width = Pt(2)

    add_text_callout(
        slide,
        Inches(0.9),
        Inches(4.45),
        Inches(11.6),
        Inches(1.85),
        "Architecture interpretation",
        "OpenCV enhances the input, YOLO localizes the object, the ROI is classified by both classical and hybrid quantum models, and OCR extracts text from the same detected object. This architecture allows a fair classifier comparison without changing the front-end computer-vision flow.",
        fill_color=WHITE,
    )


def add_workflow_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Project Presentation",
        "Experimental workflow and user interface",
        "The UI turns the project into a repeatable benchmark environment rather than a one-time code demo.",
        5,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "UI workflow",
        [
            "Train detector and ROI classifiers directly from the interface.",
            "Reuse saved artifacts instead of retraining models repeatedly.",
            "Generate benchmark reports and compare all supported models.",
            "Run full-scene inference and ROI-only inference from the same application.",
        ],
        title_color=CLASSICAL,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Research value",
        [
            "The UI standardizes training, benchmarking, and inference for fairer evaluation.",
            "Charts, tables, and saved outputs make the results presentation-ready.",
            "The workflow supports both engineering demonstration and QML experimentation.",
        ],
        title_color=QUANTUM,
        fill_color=SOFT_PURPLE,
        font_size=16,
    )


def add_dataset_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Project Presentation",
        "Dataset design and evaluation setup",
        "Detection and classification were trained on separate datasets so each research question could be measured correctly.",
        6,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.3),
        "Detection dataset",
        [
            "YOLO-style full-scene images with train and validation splits.",
            "Bounding boxes stored in standard YOLO text format.",
            "Supports realistic scene-level testing for localization and cropping.",
        ],
        title_color=CLASSICAL,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.3),
        "ROI classifier dataset",
        [
            "Cropped object images organized by class.",
            "Supports controlled feature extraction and fair classifier comparison.",
            "Balanced training and augmentation were used before model evaluation.",
        ],
        title_color=QUANTUM,
        fill_color=WHITE,
        font_size=16,
    )
    add_metric_card(slide, Inches(1.0), Inches(6.0), "ROI Samples", "997", "Classifier dataset size", color=TEAL)
    add_metric_card(slide, Inches(3.45), Inches(6.0), "Test Samples", "200", "Held-out ROI benchmark", color=GOLD)
    add_metric_card(slide, Inches(5.9), Inches(6.0), "Balanced Train", "1047", "Post-augmentation", color=CLASSICAL)
    add_metric_card(slide, Inches(8.35), Inches(6.0), "Objects", "3", "Benchmark object categories", color=QUANTUM)


def add_methods_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Methods and Technology",
        "Pipeline methods and implementation blocks",
        "The project combines standard vision methods with QML-specific classifier experiments.",
        7,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "Vision and OCR methods",
        [
            "Noise reduction, contrast enhancement, and binarization using OpenCV.",
            "YOLO-based detector for object localization.",
            "ROI extraction with controlled padding after detection.",
            "OCR ensemble using Tesseract and TrOCR.",
        ],
        title_color=CLASSICAL,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Classifier methods",
        [
            "Classical feature extraction from ROI images.",
            "Classical baselines: SVM, Logistic Regression, Random Forest, and MLP.",
            "Quantum-oriented models: QSVM variants and VQC variants implemented with Qiskit.",
            "All models were evaluated under the same held-out ROI benchmark.",
        ],
        title_color=QUANTUM,
        fill_color=WHITE,
        font_size=16,
    )


def add_qml_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Methods and Technology",
        "Hybrid QML design and feature encoding",
        "Hybrid quantum design was selected because direct quantum image processing is not practical for this object-detection research setting today.",
        8,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "Why hybrid quantum",
        [
            "Raw images are too high-dimensional for a practical pure-quantum pipeline in the current setting.",
            "Hybrid QML keeps preprocessing and feature extraction classical, then evaluates quantum models at the classifier stage.",
            "This allows quantum methods to be studied without weakening the engineering foundation of the pipeline.",
        ],
        title_color=TEAL,
        fill_color=SOFT_GREEN,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Quantum encoding and models",
        [
            "Angle encoding maps classical features to quantum rotation parameters.",
            "Amplitude encoding normalizes features into a quantum-state representation.",
            "QSVMs test quantum kernels, while VQCs test trainable variational circuits.",
            "This design supports both accuracy comparison and methodological analysis.",
        ],
        title_color=QUANTUM,
        fill_color=SOFT_PURPLE,
        font_size=16,
    )


def add_comparison_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Methods and Technology",
        "Comparison framework: Classical, hybrid quantum, and pure quantum",
        "The project formalizes each approach so the recommendation can be explained clearly to technical and non-technical audiences.",
        9,
        total,
    )
    table = slide.shapes.add_table(4, 5, Inches(0.55), Inches(2.0), Inches(12.25), Inches(4.7)).table
    headers = ["Approach", "How It Works Here", "Advantages", "Limitations", "Best Use"]
    rows = [
        ["Classical", "OpenCV + ROI features + classical classifier", "Strong baseline, faster training, easier deployment", "Lower research novelty", "Immediate production"],
        ["Hybrid Quantum", "Classical preprocessing + compressed features + quantum classifier", "Most realistic way to evaluate QML today", "Higher complexity and longer training time", "Primary research track"],
        ["Pure Quantum", "Direct quantum image encoding and classification", "High novelty and long-term strategic value", "Not practical for this task today", "Future exploration"],
    ]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = value
    style_table(table)


def add_technology_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Methods and Technology",
        "Technologies used and key skills developed",
        "The internship required integration across computer vision, QML, OCR, UI development, and technical reporting.",
        10,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "Technology stack",
        [
            "Python, OpenCV, PyTorch, and Ultralytics YOLO",
            "Qiskit and Qiskit Machine Learning",
            "Tesseract OCR and Microsoft TrOCR",
            "Flask, React, JSON-based reporting, and PowerPoint automation",
        ],
        title_color=NAVY,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Skills developed",
        [
            "End-to-end AI pipeline design and integration",
            "Dataset preparation and evaluation methodology",
            "QML benchmarking, interpretation, and recommendation",
            "Presentation-ready documentation and visual communication",
        ],
        title_color=TEAL,
        fill_color=SOFT_BLUE,
        font_size=16,
    )


def add_goals_slide_one(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Internship Objectives",
        "Achievement of internship objectives I",
        "This slide covers application, knowledge, exploration, and identification.",
        11,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "Application and knowledge",
        [
            "Applied computer vision, OCR, and machine-learning concepts in a complete working pipeline.",
            "Learned how quantum feature encoding and QML classifiers fit into a realistic hybrid workflow.",
            "Converted theoretical QML concepts into measurable implementation outcomes.",
        ],
        title_color=CLASSICAL,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Exploration and identification",
        [
            "Explored multiple encodings, kernels, and classifier families.",
            "Identified where quantum models are competitive and where classical models remain stronger.",
            "Interpreted the comparison using held-out metrics instead of assumptions.",
        ],
        title_color=QUANTUM,
        fill_color=SOFT_PURPLE,
        font_size=16,
    )


def add_goals_slide_two(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Internship Objectives",
        "Achievement of internship objectives II",
        "This slide covers innovation, engagement, evaluation, and demonstration.",
        12,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.9),
        Inches(4.7),
        "Innovation and engagement",
        [
            "Built a hybrid classical-quantum benchmark platform rather than a single-model demo.",
            "Worked across vision, OCR, QML, UI, benchmarking, and technical documentation.",
            "Created a project that supports both engineering use and research communication.",
        ],
        title_color=TEAL,
        fill_color=SOFT_GREEN,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.85),
        Inches(4.7),
        "Evaluation and demonstration",
        [
            "Benchmarked all classifier families on the same ROI dataset and evaluation logic.",
            "Generated reusable artifacts, charts, and tables for repeatable review.",
            "Prepared a presentation-ready UI and formal reporting material for evaluation.",
        ],
        title_color=GOLD,
        fill_color=WHITE,
        font_size=16,
    )


def add_results_slide_one(prs: Presentation, detector_summary: dict, benchmark: dict, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Results",
        "Detector performance and all-model accuracy comparison",
        "This slide shows the detector outcome and the ROI classification accuracy of every available model.",
        13,
        total,
    )

    add_metric_card(slide, Inches(0.75), Inches(2.0), "Precision", f"{detector_summary.get('precision', 0.0):.4f}", "Detector metric", color=CLASSICAL)
    add_metric_card(slide, Inches(3.15), Inches(2.0), "Recall", f"{detector_summary.get('recall', 0.0):.4f}", "Detector metric", color=TEAL)
    add_metric_card(slide, Inches(5.55), Inches(2.0), "mAP@50", f"{detector_summary.get('mAP50', 0.0):.4f}", "Detector metric", color=GOLD)
    add_metric_card(slide, Inches(7.95), Inches(2.0), "mAP@50-95", f"{detector_summary.get('mAP50_95', 0.0):.4f}", "Detector metric", color=CLASSICAL)

    classifier_rows = benchmark.get("classifier_benchmarks", [])
    categories = [short_model_label(str(row.get("name", ""))) for row in classifier_rows]
    values = [float(row.get("accuracy", 0.0)) * 100 for row in classifier_rows]
    add_horizontal_bar_chart(
        slide,
        Inches(0.7),
        Inches(3.25),
        Inches(8.15),
        Inches(3.5),
        "ROI classifier accuracy across all available models",
        categories,
        values,
        color=CLASSICAL,
    )

    add_bullet_panel(
        slide,
        Inches(9.0),
        Inches(3.2),
        Inches(3.7),
        Inches(3.55),
        "Key observations",
        [
            "All classifier models currently available in the project are included in the comparison.",
            "The strongest classical models and the amplitude-encoding QSVM all reached 99.5% ROI accuracy.",
            "The broader quantum family shows higher performance variability than the classical family.",
        ],
        title_color=NAVY,
        fill_color=WHITE,
        font_size=14,
    )


def add_results_slide_two(prs: Presentation, benchmark: dict, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Results",
        "Formal benchmark table and recommendation",
        "This slide provides the detailed model comparison and the final production-versus-research recommendation.",
        14,
        total,
    )

    classifiers = benchmark.get("classifier_benchmarks", [])
    first_half = classifiers[:5]
    second_half = classifiers[5:]
    add_model_table(slide, Inches(0.55), Inches(2.0), Inches(6.0), Inches(3.15), first_half)
    add_model_table(slide, Inches(6.78), Inches(2.0), Inches(6.0), Inches(3.15), second_half)

    best_classical = get_best_by_family(classifiers, "classical")
    best_quantum = get_best_by_family(classifiers, "quantum")
    family_summary = benchmark.get("classifier_family_summary", [])
    family_map = {item.get("family"): item for item in family_summary}

    add_metric_card(
        slide,
        Inches(0.7),
        Inches(5.45),
        "Best Classical",
        short_model_label(str(best_classical.get("name", ""))).replace(" [C]", ""),
        format_pct(float(best_classical.get("accuracy", 0.0))),
        color=CLASSICAL,
    )
    add_metric_card(
        slide,
        Inches(3.15),
        Inches(5.45),
        "Best Quantum",
        short_model_label(str(best_quantum.get("name", ""))).replace(" [Q]", ""),
        format_pct(float(best_quantum.get("accuracy", 0.0))),
        color=QUANTUM,
    )
    add_metric_card(
        slide,
        Inches(5.60),
        Inches(5.45),
        "Family Average",
        f"C {family_map.get('classical', {}).get('avg_accuracy', 0.0) * 100:.2f}%",
        f"Q {family_map.get('quantum', {}).get('avg_accuracy', 0.0) * 100:.2f}%",
        color=TEAL,
    )

    recommendation = benchmark.get("recommended_pipeline", {})
    add_bullet_panel(
        slide,
        Inches(8.1),
        Inches(5.28),
        Inches(4.55),
        Inches(1.55),
        "Recommendation",
        [
            f"Production recommendation: YOLO + {recommendation.get('classifier_name', 'N/A')} + OCR ensemble",
            "Research recommendation: continue hybrid QML, especially kernel-based approaches that already show competitive accuracy.",
        ],
        title_color=NAVY,
        fill_color=SOFT_BLUE,
        font_size=12,
    )


def add_future_scope_slide(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Future Scope",
        "Conclusion and future scope",
        "The project ends with a balanced conclusion: classical models lead deployment today, while hybrid QML remains the key research direction.",
        15,
        total,
    )
    add_bullet_panel(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(6.0),
        Inches(4.7),
        "Conclusion",
        [
            "The project successfully demonstrates a complete pipeline for detection, classification, and OCR.",
            "It also provides a fair and reproducible comparison between classical and hybrid quantum models.",
            "The strongest current production path is still the classical stack, led by YOLO and Classical SVM.",
            "The strongest research outcome is that a hybrid quantum model can match top ROI accuracy in this benchmark.",
        ],
        title_color=CLASSICAL,
        fill_color=WHITE,
        font_size=16,
    )
    add_bullet_panel(
        slide,
        Inches(6.85),
        Inches(2.0),
        Inches(5.8),
        Inches(4.7),
        "Future scope",
        [
            "Add stronger end-to-end validation on full-scene test images.",
            "Improve OCR robustness on difficult object text regions and reflective labels.",
            "Explore richer quantum kernels, encodings, and feature-compression strategies.",
            "Continue QML experimentation only where measurable value is demonstrated over classical baselines.",
        ],
        title_color=QUANTUM,
        fill_color=SOFT_PURPLE,
        font_size=16,
    )


def build_presentation() -> Path:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    benchmark = load_json(BENCHMARK_PATH)
    detector_summary = load_json(DETECTOR_SUMMARY_PATH)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total_slides = 15

    add_title_slide(prs, total_slides)
    add_agenda_slide(prs, total_slides)
    add_project_context_slide(prs, total_slides)
    add_architecture_slide(prs, total_slides)
    add_workflow_slide(prs, total_slides)
    add_dataset_slide(prs, total_slides)
    add_methods_slide(prs, total_slides)
    add_qml_slide(prs, total_slides)
    add_comparison_slide(prs, total_slides)
    add_technology_slide(prs, total_slides)
    add_goals_slide_one(prs, total_slides)
    add_goals_slide_two(prs, total_slides)
    add_results_slide_one(prs, detector_summary, benchmark, total_slides)
    add_results_slide_two(prs, benchmark, total_slides)
    add_future_scope_slide(prs, total_slides)

    try:
        prs.save(OUTPUT_PPTX)
        return OUTPUT_PPTX
    except PermissionError:
        try:
            prs.save(FALLBACK_OUTPUT_PPTX)
            return FALLBACK_OUTPUT_PPTX
        except PermissionError:
            timestamped = DOCS_DIR / f"INTERNSHIP_EVALUATION_PRESENTATION_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(timestamped)
            return timestamped


if __name__ == "__main__":
    output = build_presentation()
    print(f"Saved PowerPoint to: {output}")
